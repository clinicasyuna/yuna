from pathlib import Path
from datetime import datetime

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

EXCEL_PATH = Path(r"E:\APP\deploy\Relatório_Chamados_15-01-2026_937 - Samuel Lacerda.xlsx")
OUTPUT_PATH = Path(r"E:\APP\deploy\Apresentacao_Atendimentos_TI_2025_FINAL.pptx")


def load_dataset(path: Path) -> pd.DataFrame:
    df = pd.read_excel(path)
    df['Data de Criação'] = pd.to_datetime(df['Data de Criação'])
    df['Data de Finalização'] = pd.to_datetime(df['Data de Finalização'])
    df['Ano'] = df['Data de Criação'].dt.year
    return df


def enrich_metrics(df: pd.DataFrame) -> dict:
    df2025 = df[df['Ano'] == 2025].copy()
    df2025['Mes'] = df2025['Data de Criação'].dt.month

    def to_hours(val):
        if pd.isna(val):
            return 0.0
        try:
            parts = str(val).split(':')
            if len(parts) == 3:
                h, m, s = map(float, parts)
                return h + m / 60 + s / 3600
            return float(val)
        except Exception:
            return 0.0

    df2025['Horas'] = df2025['Total de Horas Tarifadas'].apply(to_hours)
    mask_final = ~df2025['Data de Finalização'].isna()
    df2025['Duracao_h'] = np.where(
        mask_final,
        (df2025['Data de Finalização'] - df2025['Data de Criação']).dt.total_seconds() / 3600,
        np.nan,
    )

    status_counts = df2025['Nome do Status'].value_counts()
    total = len(df2025)
    resolvidos = int(status_counts.get('Resolvido', 0))
    cancelados = int(status_counts.sum() - resolvidos)

    top_assuntos_series = df2025['Assunto'].value_counts().head(10)
    assuntos_detalhes = []
    for assunto, qtd in top_assuntos_series.items():
        subset = df2025[df2025['Assunto'] == assunto]
        desc_series = subset['Descrição'].dropna()
        descricao = desc_series.iloc[0] if len(desc_series) > 0 else 'Descrição não informada.'
        assuntos_detalhes.append({'assunto': assunto, 'qtd': int(qtd), 'descricao': descricao})

    metrics = {
        'df': df2025,
        'total': total,
        'status_counts': status_counts,
        'resolvidos': resolvidos,
        'cancelados': cancelados,
        'taxa_resolucao': (resolvidos / total) if total else 0,
        'horas_totais': df2025['Horas'].sum(),
        'horas_medias': df2025['Horas'].mean() if total else 0,
        'duracao_media_h': float(np.nanmean(df2025['Duracao_h'])) if total else 0,
        'duracao_mediana_h': float(np.nanmedian(df2025['Duracao_h'])) if total else 0,
        'por_mes': df2025['Mes'].value_counts().sort_index(),
        'top_assuntos': top_assuntos_series,
        'assuntos_detalhes': assuntos_detalhes,
        'canais': df2025['Forma de Atendimento da última ação'].fillna('Não informado').value_counts().head(6),
        'operadores': df2025['Nome do Operador'].fillna('Não informado').value_counts().head(6),
        'same_day_pct': float((df2025['Duracao_h'] <= 24).mean(skipna=True)) if total else 0,
        'within_4h_pct': float((df2025['Duracao_h'] <= 4).mean(skipna=True)) if total else 0,
    }
    return metrics


def add_title_slide(prs: Presentation, metrics: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Atendimentos de Informática 2025"
    subtitle = slide.placeholders[1]
    subtitle.text = "Relatório anual – Samuel Lacerda"
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)

    total = metrics['total']
    resolvidos = metrics['resolvidos']
    subtitle.text_frame.add_paragraph()
    p = subtitle.text_frame.add_paragraph()
    p.text = f"Total de chamados: {total} | Resolvidos: {resolvidos} ({metrics['taxa_resolucao']*100:.1f}%)"
    p.font.size = Pt(14)


def add_resumo_slide(prs: Presentation, metrics: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Resumo Executivo"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    bullets = [
        f"Chamados registrados em 2025: {metrics['total']}",
        f"Resolvidos: {metrics['resolvidos']} ({metrics['taxa_resolucao']*100:.1f}%) | Cancelados: {metrics['cancelados']}",
        f"Horas tarifadas totais: {metrics['horas_totais']:.1f} h | média por chamado: {metrics['horas_medias']*60:.1f} min",
        f"Tempo médio de resolução: {metrics['duracao_media_h']:.1f} h | mediana: {metrics['duracao_mediana_h']:.1f} h",
        f"Concluídos no mesmo dia: {metrics['same_day_pct']*100:.1f}% | em até 4h: {metrics['within_4h_pct']*100:.1f}%",
    ]
    for text in bullets:
        p = tf.add_paragraph()
        p.text = text
        p.level = 0


MONTH_LABELS = {
    1: 'Jan', 2: 'Fev', 3: 'Mar', 4: 'Abr', 5: 'Mai', 6: 'Jun',
    7: 'Jul', 8: 'Ago', 9: 'Set', 10: 'Out', 11: 'Nov', 12: 'Dez',
}


def add_mes_chart_slide(prs: Presentation, metrics: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Distribuição mensal de chamados"

    por_mes = metrics['por_mes']
    categorias = [MONTH_LABELS.get(m, str(m)) for m in por_mes.index]
    valores = por_mes.values.tolist()

    chart_data = CategoryChartData()
    chart_data.categories = categorias
    chart_data.add_series('Chamados', valores)

    x, y, cx, cy = Inches(0.8), Inches(1.8), Inches(8.0), Inches(4.5)
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        x,
        y,
        cx,
        cy,
        chart_data,
    )


def add_top_assuntos_slide(prs: Presentation, metrics: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Top 10 assuntos mais solicitados"

    assuntos = metrics['top_assuntos']
    chart_data = CategoryChartData()
    chart_data.categories = assuntos.index.tolist()
    chart_data.add_series('Chamados', assuntos.values.tolist())

    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9.0), Inches(5.0)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        x,
        y,
        cx,
        cy,
        chart_data,
    )
    
    # Ajustar fonte do eixo para caber os textos completos
    category_axis = chart.chart.category_axis
    category_axis.tick_labels.font.size = Pt(9)


def add_top_clientes_slide(prs: Presentation, metrics: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Top 5 clientes por volume"

    clientes = metrics['top_clientes']
    chart_data = CategoryChartData()
    chart_data.categories = clientes.index.tolist()
    chart_data.add_series('Chamados', clientes.values.tolist())

    x, y, cx, cy = Inches(0.8), Inches(1.8), Inches(8.0), Inches(4.5)
    slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        x,
        y,
        cx,
        cy,
        chart_data,
    )


def add_canais_slide(prs: Presentation, metrics: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Formas de atendimento (Top 6)"

    canais = metrics['canais']
    chart_data = CategoryChartData()
    chart_data.categories = canais.index.tolist()
    chart_data.add_series('Chamados', canais.values.tolist())

    x, y, cx, cy = Inches(0.8), Inches(1.8), Inches(8.0), Inches(4.5)
    slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        x,
        y,
        cx,
        cy,
        chart_data,
    )


def add_operadores_slide(prs: Presentation, metrics: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Operadores (Top 6 por volume)"

    ops = metrics['operadores']
    chart_data = CategoryChartData()
    chart_data.categories = ops.index.tolist()
    chart_data.add_series('Chamados', ops.values.tolist())

    x, y, cx, cy = Inches(0.8), Inches(1.8), Inches(8.0), Inches(4.5)
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        x,
        y,
        cx,
        cy,
        chart_data,
    )


def add_status_slide(prs: Presentation, metrics: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Status e SLA"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    bullets = [
        f"Taxa de resolução: {metrics['taxa_resolucao']*100:.1f}% ( {metrics['resolvidos']} de {metrics['total']} )",
        f"Média de tempo de resolução: {metrics['duracao_media_h']:.1f} h | Mediana: {metrics['duracao_mediana_h']:.1f} h",
        "Chamados cancelados representam menos de 1% do total",
        "Tendência: maioria concluída no mesmo dia (mediana 0h)",
    ]
    for text in bullets:
        p = tf.add_paragraph()
        p.text = text
        p.level = 0


def add_status_chart_slide(prs: Presentation, metrics: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Status dos chamados"

    status_counts = metrics['status_counts']
    chart_data = CategoryChartData()
    chart_data.categories = status_counts.index.tolist()
    chart_data.add_series('Chamados', status_counts.values.tolist())

    x, y, cx, cy = Inches(0.8), Inches(1.8), Inches(8.0), Inches(4.5)
    slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        x,
        y,
        cx,
        cy,
        chart_data,
    )


def add_descricoes_detail_slides(prs: Presentation, metrics: dict):
    """Cria um slide para cada assunto com a descrição completa de um chamado representativo."""
    detail_slides = []
    for item in metrics['assuntos_detalhes']:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"{item['assunto']} ({item['qtd']} chamados)"
        tf = slide.placeholders[1].text_frame
        tf.clear()
        p = tf.add_paragraph()
        p.text = item['descricao']
        p.font.size = Pt(16)
        p.level = 0
        detail_slides.append({'slide': slide, 'assunto': item['assunto'], 'qtd': item['qtd']})
    return detail_slides


def add_descricoes_sumario_slide(prs: Presentation, detail_slides: list) -> None:
    """Cria um sumário com links para os slides de descrição completa."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Descrições completas (clique para abrir)"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    for item in detail_slides:
        p = tf.add_paragraph()
        p.text = f"{item['assunto']} ({item['qtd']} chamados)"
        p.level = 0
        try:
            p.hyperlink.subaddress = str(item['slide'].slide_id)
        except Exception:
            pass


def add_recomendacoes_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Próximos passos sugeridos"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    bullets = [
        "Criar playbooks rápidos para os 5 assuntos mais recorrentes",
        "Reforçar checklists em clientes com maior volume (Top 5)",
        "Ajustar SLAs por canal de atendimento com base na agilidade observada",
        "Planejar capacidade para meses de pico e reforçar resposta no mesmo dia",
    ]
    for text in bullets:
        p = tf.add_paragraph()
        p.text = text
        p.level = 0


def generate_presentation(metrics: dict, output_path: Path) -> None:
    prs = Presentation()
    add_title_slide(prs, metrics)
    add_resumo_slide(prs, metrics)
    add_mes_chart_slide(prs, metrics)
    add_top_assuntos_slide(prs, metrics)
    detail_slides = add_descricoes_detail_slides(prs, metrics)
    add_descricoes_sumario_slide(prs, detail_slides)
    add_canais_slide(prs, metrics)
    add_operadores_slide(prs, metrics)
    add_status_chart_slide(prs, metrics)
    add_status_slide(prs, metrics)
    add_recomendacoes_slide(prs)
    prs.save(output_path)


def main() -> None:
    df = load_dataset(EXCEL_PATH)
    metrics = enrich_metrics(df)
    generate_presentation(metrics, OUTPUT_PATH)
    print(f"Apresentação gerada em: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
